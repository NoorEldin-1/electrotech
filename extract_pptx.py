import sys
import subprocess

def install_and_import():
    try:
        import pptx
    except ImportError:
        print("Installing python-pptx...")
        subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])

install_and_import()
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def extract_text(shape):
    text = ""
    if hasattr(shape, "text") and shape.text.strip():
        text += shape.text.strip() + "\n\n"
    if shape.has_table:
        table = shape.table
        for row in table.rows:
            row_data = []
            for cell in row.cells:
                row_data.append(cell.text.replace("\n", " ").strip())
            text += "| " + " | ".join(row_data) + " |\n"
            if row == table.rows[0]:
                text += "|" + "|".join(["---"] * len(row.cells)) + "|\n"
        text += "\n"
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child_shape in shape.shapes:
            text += extract_text(child_shape)
    return text

def extract_text_from_pptx(pptx_path, md_path):
    prs = Presentation(pptx_path)
    with open(md_path, 'w', encoding='utf-8') as f:
        f.write("# محتوى العرض التقديمي (المشتريات)\n\n")
        for i, slide in enumerate(prs.slides):
            f.write(f"## شريحة {i+1}\n\n")
            
            # Extract notes if any
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    f.write(f"**ملاحظات الشريحة:**\n{notes}\n\n")
                    
            for shape in slide.shapes:
                text = extract_text(shape)
                if text:
                    f.write(text)
            f.write("---\n\n")
    print("Successfully extracted text.")

if __name__ == "__main__":
    if len(sys.argv) != 3:
        print("Usage: python extract_pptx.py <input.pptx> <output.md>")
        sys.exit(1)
    extract_text_from_pptx(sys.argv[1], sys.argv[2])
